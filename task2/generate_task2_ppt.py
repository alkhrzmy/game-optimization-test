"""
Task 2 - Presentation generator (max 5 slides).
Mirrors the deck style of Task 1.
Run:  python generate_task2_ppt.py
Out:  Task2_Data_Enhancement_Gemini.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG       = RGBColor(0x10, 0x14, 0x1C)
PANEL    = RGBColor(0x18, 0x20, 0x2E)
PANEL_2  = RGBColor(0x1F, 0x29, 0x3B)
ACCENT   = RGBColor(0xFF, 0xB3, 0x00)
ACCENT_2 = RGBColor(0x35, 0xC4, 0xFF)
TEXT     = RGBColor(0xEC, 0xF0, 0xF6)
MUTED    = RGBColor(0x9A, 0xA6, 0xB8)
GOOD     = RGBColor(0x49, 0xD6, 0x8B)
WARN     = RGBColor(0xFF, 0x6B, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _bg(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG


def rect(s, x, y, w, h, fill=PANEL, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def run(p, text, size=16, color=TEXT, bold=False, italic=False, font="Segoe UI"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def para(tf, text, size=16, color=TEXT, bold=False, italic=False, after=6, before=0,
         level=0, align=PP_ALIGN.LEFT, bullet=False, first=False, font="Segoe UI"):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    run(p, ("•  " if bullet else "") + text, size=size, color=color, bold=bold,
        italic=italic, font=font)
    return p


def slide(kicker, title):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    rect(s, 0, 0, Inches(0.22), Inches(7.5), fill=ACCENT)
    tf = tbox(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.95))
    para(tf, kicker, size=13, color=ACCENT_2, bold=True, first=True, after=2)
    para(tf, title, size=29, color=TEXT, bold=True)
    rect(s, Inches(0.6), Inches(1.2), Inches(2.2), Pt(4), fill=ACCENT)
    return s


# --------------------------------------------------------------------------- #
# Slide 1 — Intro
# --------------------------------------------------------------------------- #
def s1_intro():
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    rect(s, 0, 0, Inches(0.22), Inches(7.5), fill=ACCENT)
    tf = tbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.4))
    para(tf, "TASK 2", size=15, color=ACCENT_2, bold=True, first=True, after=8)
    para(tf, "Video Game Data Enhancement", size=40, color=TEXT, bold=True, after=4)
    para(tf, "Enriching a game dataset with Google AI Studio (Gemini) API", size=20, color=MUTED)

    rect(s, Inches(0.9), Inches(4.4), Inches(5.6), Inches(2.3), fill=PANEL)
    tf = tbox(s, Inches(1.15), Inches(4.55), Inches(5.1), Inches(2.0))
    para(tf, "Goal", size=16, color=ACCENT, bold=True, first=True, after=6)
    para(tf, "Add 3 AI-generated columns to 50 game titles:", size=14, color=TEXT, after=6)
    para(tf, "genre  →  single word", size=13.5, color=GOOD, bullet=True, after=3)
    para(tf, "short_description  →  < 30 words", size=13.5, color=GOOD, bullet=True, after=3)
    para(tf, "player_mode  →  Single / Multi / Both", size=13.5, color=GOOD, bullet=True)

    rect(s, Inches(6.8), Inches(4.4), Inches(5.6), Inches(2.3), fill=PANEL_2)
    tf = tbox(s, Inches(7.05), Inches(4.55), Inches(5.1), Inches(2.0))
    para(tf, "Tools", size=16, color=ACCENT, bold=True, first=True, after=6)
    para(tf, "Python 3 + pandas (data wrangling)", size=13.5, color=TEXT, bullet=True, after=4)
    para(tf, "google-genai SDK  →  Gemini 3.1 Flash-Lite", size=13.5, color=TEXT, bullet=True, after=4)
    para(tf, "python-dotenv  →  secure API key (.env)", size=13.5, color=TEXT, bullet=True, after=4)
    para(tf, "Structured JSON output + validation layer", size=13.5, color=TEXT, bullet=True)


# --------------------------------------------------------------------------- #
# Slide 2 — Prompt Engineering: good vs bad
# --------------------------------------------------------------------------- #
def s2_prompt_good_bad():
    s = slide("PROMPT ENGINEERING (1/2)", "Good vs Bad Prompts")

    rect(s, Inches(0.6), Inches(1.5), Inches(5.95), Inches(5.3), fill=PANEL)
    rect(s, Inches(0.6), Inches(1.5), Inches(5.95), Pt(5), fill=WARN)
    tf = tbox(s, Inches(0.85), Inches(1.65), Inches(5.45), Inches(5.0))
    para(tf, "✗  BAD PROMPT", size=16, color=WARN, bold=True, first=True, after=6)
    para(tf, '"What genre is Street Fighter 6 and tell me about it?"',
         size=13.5, color=TEXT, italic=True, after=10, font="Consolas")
    para(tf, "Why it fails", size=13, color=WARN, bold=True, after=3)
    for t in ["Returns a paragraph, not a single word.",
              "No length limit on the description.",
              "player_mode free-form ('online co-op'…) → unusable.",
              "Hard to parse → breaks the pipeline."]:
        para(tf, t, size=13, color=TEXT, bullet=True, after=5)

    rect(s, Inches(6.78), Inches(1.5), Inches(5.95), Inches(5.3), fill=PANEL)
    rect(s, Inches(6.78), Inches(1.5), Inches(5.95), Pt(5), fill=GOOD)
    tf = tbox(s, Inches(7.03), Inches(1.65), Inches(5.45), Inches(5.0))
    para(tf, "✓  GOOD PROMPT", size=16, color=GOOD, bold=True, first=True, after=6)
    para(tf, 'Role + strict JSON schema + explicit rules + temp=0',
         size=13.5, color=TEXT, italic=True, after=10, font="Consolas")
    para(tf, "Why it works", size=13, color=GOOD, bold=True, after=3)
    for t in ['Assigns a role: "precise metadata classifier".',
              'Demands ONE word for genre, < 30 words for description.',
              'Constrains player_mode to 3 exact enum values.',
              'Forces raw JSON only → reliable machine parsing.',
              'temperature = 0 → deterministic, repeatable.']:
        para(tf, t, size=13, color=TEXT, bullet=True, after=5)


# --------------------------------------------------------------------------- #
# Slide 3 — Prompt Engineering: the final prompt + sample output
# --------------------------------------------------------------------------- #
def s3_prompt_final():
    s = slide("PROMPT ENGINEERING (2/2)", "Final Prompt & Sample Output")

    rect(s, Inches(0.6), Inches(1.5), Inches(6.4), Inches(5.3), fill=PANEL)
    tf = tbox(s, Inches(0.82), Inches(1.62), Inches(5.95), Inches(5.0))
    para(tf, "The production prompt (per title)", size=15, color=ACCENT, bold=True, first=True, after=6)
    code = [
        'You are a precise video-game',
        'metadata classifier. For "{title}"',
        'return ONLY JSON with keys:',
        '  genre            (ONE word)',
        '  short_description(< 30 words)',
        '  player_mode      (Single/Multi/Both)',
        'Rules: single-word genre; <30 words;',
        'exact enum; best guess, never refuse;',
        'raw JSON only, no code fences.',
    ]
    for line in code:
        para(tf, line, size=13, color=TEXT, after=2, font="Consolas")
    para(tf, "One structured call per game replaces three vague ones → fewer requests, consistent output.",
         size=12.5, color=MUTED, before=8)

    rect(s, Inches(7.2), Inches(1.5), Inches(5.53), Inches(5.3), fill=PANEL_2)
    tf = tbox(s, Inches(7.42), Inches(1.62), Inches(5.1), Inches(5.0))
    para(tf, "Sample output  (Street Fighter 6)", size=15, color=ACCENT, bold=True, first=True, after=6)
    out = [
        '{',
        '  "genre": "Fighting",',
        '  "short_description": "Competitive',
        '   1v1 fighting game with',
        '   real-time combos and special',
        '   moves across iconic characters.",',
        '  "player_mode": "Both"',
        '}',
    ]
    for line in out:
        para(tf, line, size=13, color=GOOD, after=2, font="Consolas")
    para(tf, "Post-processing then enforces the rules: first-word genre, "
             "word-count cap, enum mapping.", size=12.5, color=MUTED, before=8)


# --------------------------------------------------------------------------- #
# Slide 4 — Challenges & Solutions
# --------------------------------------------------------------------------- #
def s4_challenges():
    s = slide("CHALLENGES & SOLUTIONS", "What Went Wrong & How It Was Fixed")
    rows = [
        ("Rate limits (HTTP 429) on free tier",
         "Exponential backoff (2^n) + configurable --sleep between calls."),
        ("Model adds prose / ```json fences",
         "JSON extractor strips fences and grabs the outer {...} before parsing."),
        ("Multi-word or fuzzy genre output",
         "clean_genre() keeps only the first token; normalises to one word."),
        ("Descriptions exceeding 30 words",
         "clean_description() hard-caps to 29 words + ellipsis."),
        ("Inconsistent player_mode wording",
         "clean_player_mode() maps any variant onto the 3 allowed enum labels."),
        ("Crash mid-run wastes quota",
         "Incremental save + resumable: re-run skips already-enriched rows."),
    ]
    y = Inches(1.55)
    for i, (chal, sol) in enumerate(rows):
        fill = PANEL if i % 2 == 0 else PANEL_2
        rect(s, Inches(0.6), y, Inches(12.13), Inches(0.82), fill=fill)
        tf = tbox(s, Inches(0.8), y, Inches(5.6), Inches(0.82), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, chal, size=13.5, color=WARN, bold=True, first=True)
        tf = tbox(s, Inches(6.5), y, Inches(6.1), Inches(0.82), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, sol, size=13.5, color=TEXT, first=True)
        y = y + Inches(0.87)


# --------------------------------------------------------------------------- #
# Slide 5 — Conclusion & Reflections
# --------------------------------------------------------------------------- #
def s5_conclusion():
    s = slide("CONCLUSION & REFLECTIONS", "Lessons Learned & Next Steps")
    rect(s, Inches(0.6), Inches(1.55), Inches(5.95), Inches(5.25), fill=PANEL)
    tf = tbox(s, Inches(0.85), Inches(1.7), Inches(5.45), Inches(5.0))
    para(tf, "Lessons learned", size=16, color=ACCENT, bold=True, first=True, after=8)
    for t in ["Constraints in the prompt matter more than length — schema + enums = clean data.",
              "Always validate model output; never trust it raw.",
              "Determinism (temp=0) makes results reproducible and gradeable.",
              "Cheap resilience (retry + resume) saves real money on quota."]:
        para(tf, t, size=14, color=TEXT, bullet=True, after=9)

    rect(s, Inches(6.78), Inches(1.55), Inches(5.95), Inches(5.25), fill=PANEL_2)
    tf = tbox(s, Inches(7.03), Inches(1.7), Inches(5.45), Inches(5.0))
    para(tf, "Future improvements", size=16, color=ACCENT, bold=True, first=True, after=8)
    for t in ["Batch / async calls for higher throughput on large datasets.",
              "Use the image_url (multimodal) to improve genre accuracy.",
              "response_schema (typed) instead of prompt-only JSON.",
              "Confidence score + human review queue for uncertain titles.",
              "Cache by title to dedupe across dataset versions."]:
        para(tf, t, size=14, color=TEXT, bullet=True, after=9)


def main():
    s1_intro()
    s2_prompt_good_bad()
    s3_prompt_final()
    s4_challenges()
    s5_conclusion()
    out = "Task2_Data_Enhancement_Gemini.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides.")


if __name__ == "__main__":
    main()
