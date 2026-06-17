"""
Task 3 - Presentation generator (Data Analysis).
Same dark deck style as Task 1 & 2. Numbers are the real query outputs.
Run:  python generate_task3_ppt.py
Out:  Task3_Data_Analysis_Eklipse.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BG       = RGBColor(0x10, 0x14, 0x1C)
PANEL    = RGBColor(0x18, 0x20, 0x2E)
PANEL_2  = RGBColor(0x1F, 0x29, 0x3B)
ACCENT   = RGBColor(0xFF, 0xB3, 0x00)
ACCENT_2 = RGBColor(0x35, 0xC4, 0xFF)
TEXT     = RGBColor(0xEC, 0xF0, 0xF6)
MUTED    = RGBColor(0x9A, 0xA6, 0xB8)
GOOD     = RGBColor(0x49, 0xD6, 0x8B)
WARN     = RGBColor(0xFF, 0x6B, 0x6B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _bg(s):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG


def rect(s, x, y, w, h, fill=PANEL, line=None):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    sp.shadow.inherit = False
    return sp


def tbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def run(p, text, size=16, color=TEXT, bold=False, italic=False, font="Segoe UI"):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def para(tf, text, size=16, color=TEXT, bold=False, italic=False, after=6, before=0,
         level=0, align=PP_ALIGN.LEFT, bullet=False, first=False, font="Segoe UI"):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(after)
    p.space_before = Pt(before)
    run(p, ("•  " if bullet else "") + text, size=size, color=color, bold=bold,
        italic=italic, font=font)
    return p


def slide(kicker, title):
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    rect(s, 0, 0, Inches(0.22), Inches(7.5), fill=ACCENT)
    tf = tbox(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.95))
    para(tf, kicker, size=13, color=ACCENT_2, bold=True, first=True, after=2)
    para(tf, title, size=28, color=TEXT, bold=True)
    rect(s, Inches(0.6), Inches(1.2), Inches(2.2), Pt(4), fill=ACCENT)
    return s


def title_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    rect(s, 0, 0, Inches(0.22), Inches(7.5), fill=ACCENT)
    tf = tbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6))
    para(tf, "TASK 3", size=15, color=ACCENT_2, bold=True, first=True, after=8)
    para(tf, "Data Analysis", size=44, color=TEXT, bold=True, after=4)
    para(tf, "Key business metrics for Eklipse, extracted with SQL", size=20, color=MUTED)
    tf2 = tbox(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4))
    para(tf2, "5 tables  •  5 SQL queries (each joins \u22652 tables)  •  SQLite + Python",
         size=16, color=TEXT, bold=True, first=True, after=4)
    para(tf2, "Dataset window: streams, clips & premium activity (2023)", size=13, color=MUTED)


def dataset_slide():
    s = slide("STEP 1 — THE DATA", "Dataset & Approach")
    # left: tables
    rect(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.3), fill=PANEL)
    tf = tbox(s, Inches(0.85), Inches(1.65), Inches(5.5), Inches(5.0))
    para(tf, "5 tables (SQL DB export)", size=16, color=ACCENT, bold=True, first=True, after=8)
    rows = [
        ("gamesession", "37,579", "streams users submitted"),
        ("clips", "552,031", "clips generated (AI / edited)"),
        ("downloaded_clips", "12,938", "clip download actions"),
        ("shared_clips", "5,798", "clip share / schedule actions"),
        ("premium", "384", "premium subscription records"),
    ]
    for name, cnt, desc in rows:
        p = tf.add_paragraph(); p.space_after = Pt(4)
        run(p, f"{name}  ", size=14, color=ACCENT_2, bold=True, font="Consolas")
        run(p, f"{cnt}", size=14, color=GOOD, bold=True, font="Consolas")
        p2 = tf.add_paragraph(); p2.space_after = Pt(8)
        run(p2, f"   {desc}", size=12.5, color=MUTED)
    para(tf, "Linked by user_id, clips.id \u2194 clip_id, gamesession.id \u2194 gamesession_Id.",
         size=12.5, color=MUTED, before=6)

    # right: method
    rect(s, Inches(6.8), Inches(1.5), Inches(5.93), Inches(5.3), fill=PANEL_2)
    tf = tbox(s, Inches(7.05), Inches(1.65), Inches(5.45), Inches(5.0))
    para(tf, "Approach", size=16, color=ACCENT, bold=True, first=True, after=8)
    for t in ["Loaded all 5 CSVs into a SQLite database.",
              "Wrote 5 SQL queries \u2014 each JOINs \u22652 tables \u2014 to extract every metric.",
              "Ran them programmatically; used Python (pandas) only for light aggregation & the visuals.",
              "Every number in this deck is a real query output, not an estimate."]:
        para(tf, t, size=14, color=TEXT, bullet=True, after=10)


def overview_slide():
    s = slide("STEP 2 — WHAT TO TRACK", "5 Metrics Eklipse Should Watch")
    headers = ["#", "Metric", "What it answers", "Tables joined"]
    data = [
        ["1", "Clip Utilization Funnel", "Are generated clips actually used?", "clips + downloaded + shared"],
        ["2", "Free \u2192 Premium Conversion", "Is monetisation working?", "gamesession + premium"],
        ["3", "Premium Engagement Lift", "Do premium users use more?", "clips + premium"],
        ["4", "Top Games by Clip Yield", "Where is the value concentrated?", "gamesession + clips"],
        ["5", "Premium Churn vs Engagement", "Who churns, and why?", "premium + clips"],
    ]
    col_w = [Inches(0.55), Inches(3.4), Inches(4.5), Inches(3.68)]
    x0, y0 = Inches(0.6), Inches(1.55)
    rowh = Inches(0.9)
    x = x0
    for j, h in enumerate(headers):
        rect(s, x, y0, col_w[j], Inches(0.5), fill=ACCENT)
        tf = tbox(s, x, y0, col_w[j], Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, h, size=12.5, color=BG, bold=True, first=True, align=PP_ALIGN.CENTER)
        x += col_w[j]
    for i, rowd in enumerate(data):
        y = y0 + Inches(0.5) + rowh * i
        x = x0
        fill = PANEL if i % 2 == 0 else PANEL_2
        for j, cell in enumerate(rowd):
            rect(s, x, y, col_w[j], rowh, fill=fill)
            tf = tbox(s, x + Inches(0.08), y, col_w[j] - Inches(0.16), rowh, anchor=MSO_ANCHOR.MIDDLE)
            col = ACCENT_2 if j == 0 else TEXT
            mono = "Consolas" if j == 3 else "Segoe UI"
            sz = 11.5 if j == 3 else 13
            para(tf, cell, size=sz, color=col, bold=(j in (0, 1)), first=True,
                 align=PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT, font=mono)
            x += col_w[j]


def metric_slide(num, title, why, sql_tables, kpis, insight, insight_color=GOOD):
    """kpis: list of (label, value, color)."""
    s = slide(f"METRIC {num}", title)
    # why banner
    rect(s, Inches(0.6), Inches(1.4), Inches(12.13), Inches(0.95), fill=PANEL_2)
    rect(s, Inches(0.6), Inches(1.4), Pt(5), Inches(0.95), fill=ACCENT)
    tf = tbox(s, Inches(0.85), Inches(1.46), Inches(11.7), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    run(p, "Why track it:  ", size=13.5, color=ACCENT, bold=True)
    run(p, why, size=13.5, color=TEXT)

    # KPI cards
    n = len(kpis)
    gap = Inches(0.2)
    total_w = Inches(12.13)
    card_w = Emu_div(total_w, n, gap)
    x = Inches(0.6)
    for label, value, col in kpis:
        rect(s, x, Inches(2.55), card_w, Inches(1.7), fill=PANEL)
        tf = tbox(s, x, Inches(2.62), card_w, Inches(1.6), anchor=MSO_ANCHOR.MIDDLE)
        para(tf, value, size=30, color=col, bold=True, first=True, after=2, align=PP_ALIGN.CENTER)
        para(tf, label, size=12, color=MUTED, align=PP_ALIGN.CENTER)
        x = x + card_w + gap

    # SQL tables + insight
    rect(s, Inches(0.6), Inches(4.45), Inches(12.13), Inches(2.35), fill=PANEL)
    tf = tbox(s, Inches(0.85), Inches(4.58), Inches(11.7), Inches(2.1))
    p = tf.paragraphs[0]
    run(p, "SQL joins:  ", size=12.5, color=ACCENT_2, bold=True, font="Consolas")
    run(p, sql_tables, size=12.5, color=TEXT, font="Consolas")
    para(tf, "Insight & recommendation", size=13.5, color=ACCENT, bold=True, before=8, after=3)
    for line in insight:
        para(tf, line, size=13, color=TEXT, bullet=True, after=5)


def Emu_div(total, n, gap):
    # total width minus (n-1) gaps, divided by n
    from pptx.util import Emu
    return Emu(int((total - gap * (n - 1)) / n))


def summary_slide():
    s = slide("CONCLUSION", "What the Numbers Say")
    rect(s, Inches(0.6), Inches(1.5), Inches(6.0), Inches(5.3), fill=PANEL)
    tf = tbox(s, Inches(0.85), Inches(1.65), Inches(5.5), Inches(5.0))
    para(tf, "Headline findings", size=16, color=ACCENT, bold=True, first=True, after=8)
    for t in ["Clip utilization is the #1 gap: only 2.3% downloaded, 0.5% shared.",
              "Conversion is healthy at 7.7% of active users.",
              "Premium users barely out-create free users (125 vs 121) \u2014 value is in features, not volume.",
              "COD Warzone dominates volume; LoL/WoW/SMITE lead clips-per-stream.",
              "Churned premium users were less engaged (107 vs 135 clips)."]:
        para(tf, t, size=13.5, color=TEXT, bullet=True, after=8)

    rect(s, Inches(6.8), Inches(1.5), Inches(5.93), Inches(5.3), fill=PANEL_2)
    tf = tbox(s, Inches(7.05), Inches(1.65), Inches(5.45), Inches(5.0))
    para(tf, "Recommended actions", size=16, color=ACCENT, bold=True, first=True, after=8)
    for t in ["Boost clip relevance & make download/share 1-click \u2014 small lifts move a huge base.",
              "Surface premium-only value (CTT, scheduling, quality) since clip volume alone doesn't differentiate.",
              "Use early engagement (clips in first weeks) as a churn early-warning to trigger save-flows.",
              "Double down on high-yield games for model quality & partnerships.",
              "Track these 5 metrics on a weekly dashboard."]:
        para(tf, t, size=13.5, color=GOOD, bullet=True, after=8)


def main():
    title_slide()
    dataset_slide()
    overview_slide()

    metric_slide(
        1, "Clip Utilization Funnel (CTT)",
        "Eklipse's core value is clips people actually use. Generated-but-unused clips = wasted value.",
        "clips  LEFT JOIN downloaded_clips ON id=clip_id  LEFT JOIN shared_clips ON id=clip_id",
        [("Clips generated", "552,031", ACCENT_2),
         ("Downloaded", "2.34%", WARN),
         ("Shared", "0.54%", WARN)],
        ["Only ~1 in 43 generated clips is downloaded and ~1 in 185 is shared \u2014 the biggest improvement lever in the funnel.",
         "Action: improve clip relevance/ranking and reduce friction (one-click download/share to TikTok)."],
    )
    metric_slide(
        2, "Free \u2192 Premium Conversion",
        "The headline monetisation metric: what share of active users (\u22651 stream) ever go premium.",
        "gamesession  LEFT JOIN premium ON user_id",
        [("Active users", "4,322", ACCENT_2),
         ("Converted", "333", GOOD),
         ("Conversion rate", "7.7%", GOOD)],
        ["7.7% is solid for a freemium product (typical benchmark is 2\u20135%).",
         "Action: protect this funnel; A/B test upgrade prompts at high-intent moments (e.g. right after a share)."],
    )
    metric_slide(
        3, "Premium Engagement Lift",
        "If premium users create far more clips, it validates the tier; if not, value lives in features.",
        "clips (per-user) LEFT JOIN premium ON user_id",
        [("Premium avg", "125.0", GOOD),
         ("Free avg", "120.6", ACCENT_2),
         ("Lift", "+3.6%", WARN)],
        ["Premium users barely out-create free users \u2014 raw clip volume is NOT what differentiates premium.",
         "Action: lead premium messaging with feature value (CTT, scheduling, higher quality), not 'more clips'."],
    )
    metric_slide(
        4, "Top Games by Clip Yield",
        "Eklipse runs per-game AI models; volume & clips-per-stream show where to invest.",
        "gamesession  JOIN clips ON gamesession.id = clips.gamesession_Id  (GROUP BY game_name)",
        [("Top by volume", "Warzone", ACCENT_2),
         ("Warzone clips", "139,790", GOOD),
         ("Highest yield", "LoL 24.4/stream", GOOD)],
        ["COD Warzone drives the most clips overall; LoL, WoW, SMITE & Rocket League produce the most clips per stream (~24).",
         "Action: prioritise model quality for top-volume titles and feature high-yield games in growth campaigns."],
    )
    metric_slide(
        5, "Premium Churn vs Engagement",
        "Churn is the silent killer; linking it to engagement turns it into an actionable signal.",
        "premium (per-user) LEFT JOIN clips (per-user) ON user_id",
        [("Churn rate", "32%", WARN),
         ("Churned avg clips", "106.8", WARN),
         ("Retained avg clips", "134.7", GOOD)],
        ["Churned members generated ~21% fewer clips than retained ones \u2014 low engagement precedes cancellation.",
         "Action: monitor early-weeks clip activity; trigger onboarding/save-flows for low-engagement premium users."],
    )

    summary_slide()
    out = "Task3_Data_Analysis_Eklipse.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides.")


if __name__ == "__main__":
    main()
