"""
Task 1 - Game Understanding (Eklipse AI Data Optimization Test)
Generates a Google-Slides-ready .pptx deck for automated clip detection in
Garena Delta Force (Tactical FPS).

Run:
    python generate_task1_ppt.py
Output:
    Task1_Game_Understanding_DeltaForce.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Theme / palette (dark "tactical" gaming look)
# ---------------------------------------------------------------------------
BG        = RGBColor(0x10, 0x14, 0x1C)   # near-black navy
PANEL     = RGBColor(0x18, 0x20, 0x2E)   # panel background
PANEL_2   = RGBColor(0x1F, 0x29, 0x3B)   # alt panel
ACCENT    = RGBColor(0xFF, 0xB3, 0x00)   # tactical amber
ACCENT_2  = RGBColor(0x35, 0xC4, 0xFF)   # cyan
TEXT      = RGBColor(0xEC, 0xF0, 0xF6)   # near-white
MUTED     = RGBColor(0x9A, 0xA6, 0xB8)   # muted grey
GOOD      = RGBColor(0x49, 0xD6, 0x8B)
WARN      = RGBColor(0xFF, 0x6B, 0x6B)

W, H = Inches(13.333), Inches(7.5)       # 16:9 widescreen

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def _bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=PANEL, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def set_run(run, text, size=18, color=TEXT, bold=False, italic=False, font="Segoe UI"):
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font


def add_para(tf, text, size=18, color=TEXT, bold=False, italic=False,
             space_after=6, space_before=0, level=0, align=PP_ALIGN.LEFT,
             bullet=False, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    prefix = "•  " if bullet else ""
    set_run(p.add_run(), prefix + text, size=size, color=color, bold=bold, italic=italic)
    return p


def accent_bar(slide, y=Inches(1.18)):
    rect(slide, Inches(0.6), y, Inches(2.2), Pt(4), fill=ACCENT)


def slide_title(slide, kicker, title, color=ACCENT):
    _, tf = textbox(slide, Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.9))
    add_para(tf, kicker, size=13, color=color, bold=True, first=True, space_after=2)
    add_para(tf, title, size=30, color=TEXT, bold=True)
    accent_bar(slide)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def new_slide():
    s = prs.slides.add_slide(BLANK)
    _bg(s)
    return s


def title_slide():
    s = new_slide()
    # side accent
    rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
    _, tf = textbox(s, Inches(0.9), Inches(2.2), Inches(11.6), Inches(2.8))
    add_para(tf, "EKLIPSE AI DATA OPTIMIZATION TEST", size=15, color=ACCENT_2, bold=True, first=True, space_after=10)
    add_para(tf, "Task 1: Game Understanding", size=44, color=TEXT, bold=True, space_after=6)
    add_para(tf, "Automated Gameplay Clip Detection for Content Creators", size=22, color=MUTED, space_after=18)
    _, tf2 = textbox(s, Inches(0.9), Inches(5.4), Inches(11.6), Inches(1.4))
    add_para(tf2, "Game: Garena Delta Force   |   Genre: Tactical FPS (Extraction + Large-Scale Warfare)",
             size=18, color=TEXT, bold=True, first=True, space_after=4)
    add_para(tf2, "Prepared by: <Nama Anda>   |   Modes analysed: Operations (Hazard Ops) & Warfare", size=14, color=MUTED)


def game_selection_slide():
    s = new_slide()
    slide_title(s, "STEP 1 — GAME SELECTION", "Garena Delta Force — Tactical FPS")

    # Left: game profile
    rect(s, Inches(0.6), Inches(1.5), Inches(5.9), Inches(5.3), fill=PANEL)
    _, tf = textbox(s, Inches(0.85), Inches(1.65), Inches(5.4), Inches(5.0))
    add_para(tf, "Game Profile", size=18, color=ACCENT, bold=True, first=True, space_after=8)
    rows = [
        ("Title", "Garena Delta Force (Team Jade / Garena)"),
        ("Genre", "Tactical First-Person Shooter (FPS)"),
        ("Sub-genre", "Hero-shooter w/ class-based Operators"),
        ("Platforms", "PC & Mobile (cross-progression)"),
        ("Why relevant", "Fast-growing title with high clip-worthy moments for streamers"),
    ]
    for k, v in rows:
        p = add_para(tf, k, size=14, color=ACCENT_2, bold=True, space_after=0, space_before=6)
        add_para(tf, v, size=15, color=TEXT, space_after=2)

    # Right: two modes
    rect(s, Inches(6.8), Inches(1.5), Inches(5.9), Inches(2.55), fill=PANEL_2)
    _, tf1 = textbox(s, Inches(7.05), Inches(1.62), Inches(5.4), Inches(2.3))
    add_para(tf1, "MODE A — Operations (Hazard Ops)", size=16, color=ACCENT, bold=True, first=True, space_after=4)
    add_para(tf1, "Extraction shooter: drop in, loot high-value gear (incl. MandelBrick), fight AI + rival squads, and extract alive. High-stakes, gear-on-the-line tension.", size=14, color=TEXT)

    rect(s, Inches(6.8), Inches(4.25), Inches(5.9), Inches(2.55), fill=PANEL_2)
    _, tf2 = textbox(s, Inches(7.05), Inches(4.37), Inches(5.4), Inches(2.3))
    add_para(tf2, "MODE B — Warfare", size=16, color=ACCENT, bold=True, first=True, space_after=4)
    add_para(tf2, "Massive-scale multiplayer (infantry + vehicles like tanks/helis) fighting over objectives. Big team plays, vehicle dominance, and round-deciding clutches.", size=14, color=TEXT)


def methodology_slide():
    s = new_slide()
    slide_title(s, "STEP 2 — OBSERVATION & DETECTION APPROACH", "How Events Are Detected")
    _, tf = textbox(s, Inches(0.6), Inches(1.45), Inches(12.1), Inches(0.7))
    add_para(tf, "Based on 30+ minutes of gameplay footage across Operations & Warfare. Detection fuses 3 signal layers; an event only fires when ≥2 layers agree (sensor fusion) to keep false positives low.",
             size=15, color=MUTED, first=True)

    cards = [
        ("1 · UI / HUD Layer", ACCENT, [
            "Template-match HUD widgets (kill feed, timers, vehicle HUD, ability icons).",
            "OCR on text: kill feed names, 'OVERTIME', decode timer, objective banners.",
            "Watch state changes: ability ready→on-cooldown, teammate alive→dead.",
        ]),
        ("2 · Audio Layer", ACCENT_2, [
            "Kill-confirm / hitmarker stings, ability activation SFX & voice lines.",
            "Announcer cues (overtime sting), music intensity ramps.",
            "Tank cannon fire & vehicle engine ambience.",
        ]),
        ("3 · Scene / State Layer", GOOD, [
            "Scene transitions: entering vehicle, last-alive, match-result screen.",
            "Objective-state changes (capture / decode complete / extraction).",
            "Combat density (gunfire + damage indicators) spikes.",
        ]),
    ]
    x = Inches(0.6)
    for title, col, items in cards:
        rect(s, x, Inches(2.35), Inches(3.95), Inches(4.35), fill=PANEL)
        rect(s, x, Inches(2.35), Inches(3.95), Pt(5), fill=col)
        _, tf = textbox(s, x + Inches(0.22), Inches(2.55), Inches(3.5), Inches(4.0))
        add_para(tf, title, size=16, color=col, bold=True, first=True, space_after=8)
        for it in items:
            add_para(tf, it, size=13, color=TEXT, bullet=True, space_after=8)
        x = x + Inches(4.07)


def overview_slide():
    s = new_slide()
    slide_title(s, "STEP 3 — EVENT MAP", "6 Clip-Worthy Events Identified")
    from pptx.enum.shapes import MSO_SHAPE
    headers = ["#", "Event", "Mode", "Primary Signal", "Why It Matters"]
    data = [
        ["1", "MandelBrick Contest & Decoding", "Operations", "UI timer + combat", "Highest-stakes objective fight"],
        ["2", "Solo Squad Wipe", "Operations", "Kill feed burst", "Pure individual skill"],
        ["3", "Tank Multi-Kill", "Warfare", "Vehicle HUD + kill feed", "Spectacular vehicle dominance"],
        ["4", "Operator Skill Combo", "Warfare", "Ability VFX/SFX + kills", "Tactical depth & teamplay"],
        ["5", "Clutch 1v4", "Both", "Last-alive + kills", "Against-the-odds hero moment"],
        ["6", "Clutch in Overtime", "Warfare", "'OVERTIME' + decisive play", "Maximum drama / match-decider"],
    ]
    col_w = [Inches(0.6), Inches(3.7), Inches(1.7), Inches(2.9), Inches(3.2)]
    x0, y0 = Inches(0.6), Inches(1.55)
    rowh = Inches(0.78)
    # header
    x = x0
    for j, htext in enumerate(headers):
        c = rect(s, x, y0, col_w[j], Inches(0.55), fill=ACCENT)
        _, tf = textbox(s, x, y0, col_w[j], Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
        add_para(tf, htext, size=13, color=BG, bold=True, first=True, align=PP_ALIGN.CENTER)
        x += col_w[j]
    # rows
    for i, row in enumerate(data):
        y = y0 + Inches(0.55) + rowh * i
        x = x0
        fill = PANEL if i % 2 == 0 else PANEL_2
        for j, cell in enumerate(row):
            rect(s, x, y, col_w[j], rowh, fill=fill)
            _, tf = textbox(s, x + Inches(0.06), y, col_w[j] - Inches(0.12), rowh, anchor=MSO_ANCHOR.MIDDLE)
            col = ACCENT_2 if j == 0 else (TEXT if j != 1 else TEXT)
            add_para(tf, cell, size=12.5, color=col, bold=(j in (0, 1)), first=True,
                     align=PP_ALIGN.CENTER if j in (0, 2) else PP_ALIGN.LEFT)
            x += col_w[j]


def event_slide(idx, name, mode, justification, criteria, clip):
    """criteria: dict with keys ui, nonui, trigger, thresholds, fp.
       clip: dict with keys preroll, postroll, duration, naming."""
    s = new_slide()
    slide_title(s, f"EVENT {idx} · {mode.upper()}", name)

    # Justification banner
    rect(s, Inches(0.6), Inches(1.45), Inches(12.13), Inches(0.85), fill=PANEL_2)
    rect(s, Inches(0.6), Inches(1.45), Pt(5), Inches(0.85), fill=ACCENT)
    _, tf = textbox(s, Inches(0.85), Inches(1.5), Inches(11.7), Inches(0.78), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    set_run(p.add_run(), "Why capture it:  ", size=14, color=ACCENT, bold=True)
    set_run(p.add_run(), justification, size=14, color=TEXT)

    # Left column: detection criteria
    rect(s, Inches(0.6), Inches(2.45), Inches(7.55), Inches(4.45), fill=PANEL)
    _, tf = textbox(s, Inches(0.85), Inches(2.58), Inches(7.05), Inches(4.25))
    add_para(tf, "DETECTION CRITERIA", size=15, color=ACCENT, bold=True, first=True, space_after=6)
    blocks = [
        ("UI Elements", criteria["ui"], ACCENT_2),
        ("Non-UI Method", criteria["nonui"], ACCENT_2),
        ("Trigger", criteria["trigger"], GOOD),
        ("Thresholds", criteria["thresholds"], GOOD),
        ("False-Positive Mitigation", criteria["fp"], WARN),
    ]
    for label, val, col in blocks:
        p = add_para(tf, label, size=12.5, color=col, bold=True, space_after=0, space_before=4)
        add_para(tf, val, size=12.5, color=TEXT, space_after=2)

    # Right column: clip definition & management
    rect(s, Inches(8.35), Inches(2.45), Inches(4.38), Inches(4.45), fill=PANEL_2)
    _, tf = textbox(s, Inches(8.6), Inches(2.58), Inches(3.95), Inches(4.25))
    add_para(tf, "CLIP DEFINITION & MGMT", size=15, color=ACCENT, bold=True, first=True, space_after=8)
    items = [
        ("Pre-roll", clip["preroll"]),
        ("Post-roll", clip["postroll"]),
        ("Duration", clip["duration"]),
        ("Naming", clip["naming"]),
    ]
    for k, v in items:
        add_para(tf, k, size=12.5, color=GOOD, bold=True, space_after=0, space_before=4)
        add_para(tf, v, size=12.5, color=TEXT, space_after=2)


EVENTS = [
    dict(
        idx=1, name="MandelBrick Contest & Decoding", mode="Operations",
        justification="The single highest-stakes objective in extraction mode — carrying the MandelBrick exposes your position to ALL squads, turning the decode station into a multi-team firefight. Shows strategy, risk management, and clutch combat.",
        criteria=dict(
            ui="Objective banner 'MandelBrick acquired'; decode progress bar/timer at the Decoding Station; minimap decode marker; 'position exposed to all squads' warning.",
            nonui="Distinct decode start/complete audio cue; sharp rise in nearby gunfire density.",
            trigger="Decode progress UI appears (template + OCR 'decoding') AND combat activity (kill-feed / damage) occurs within the protection window.",
            thresholds="~4-min decode protection window; require ≥2 player kills OR sustained enemy contact to qualify as a 'contested' decode (skip uncontested decodes).",
            fp="Require BOTH decode-bar template AND combat signal; template-match the specific decode bar to avoid confusion with revive/capture bars.",
        ),
        clip=dict(preroll="8 s (lead-up to securing the brick)", postroll="5 s after decode completes / extraction starts",
                  duration="45–60 s (cap)", naming="DF_OPS_MandelBrick_<map>_<UTC>_<user>.mp4"),
    ),
    dict(
        idx=2, name="Solo Squad Wipe", mode="Operations",
        justification="One player eliminating an entire enemy squad single-handedly — the clearest display of individual skill and a staple of highlight reels.",
        criteria=dict(
            ui="Rapid stacked kill-feed entries attributed to the local player; 'Enemy eliminated' pop-ups; squad-wipe banner (if shown).",
            nonui="Stacking kill-confirm / hitmarker audio cues in quick succession.",
            trigger="N kills by the LOCAL player within window T, against members of one enemy squad (e.g., ≥3 kills in ≤30 s).",
            thresholds="Kills ≥3 within ≤30–45 s; killer name = local player; player still alive shortly after.",
            fp="OCR-confirm killer = local player; weight real-player kills over AI; require time-clustering & post-wipe survival.",
        ),
        clip=dict(preroll="10 s", postroll="5 s", duration="25–40 s",
                  naming="DF_OPS_SquadWipe_<map>_<UTC>_<user>.mp4"),
    ),
    dict(
        idx=3, name="Tank Multi-Kill", mode="Warfare",
        justification="Vehicle dominance on a massive map — visually spectacular and shows map control and impact. Prime Warfare highlight material.",
        criteria=dict(
            ui="Active tank/vehicle HUD overlay (vehicle ammo, reticle, vehicle health) + kill feed showing the tank-cannon weapon icon on multiple kills.",
            nonui="Tank main-gun fire audio; vehicle engine ambience confirming the player is the gunner/driver.",
            trigger="Tank HUD active AND ≥2–3 kills within ~10–15 s attributed to the player via the tank weapon icon.",
            thresholds="Multi-kill ≥2 (double) / ≥3 (triple) in window; weapon icon = tank main gun (not infantry weapon).",
            fp="Confirm tank-HUD template (distinct from infantry HUD); verify kill-feed weapon icon = tank cannon; tag roadkills separately.",
        ),
        clip=dict(preroll="6 s", postroll="4 s", duration="20–35 s",
                  naming="DF_WAR_TankMultiKill_<map>_<UTC>_<user>.mp4"),
    ),
    dict(
        idx=4, name="Special Operator Skill Combination", mode="Warfare",
        justification="Showcases the game's tactical depth — chaining Operator abilities/ultimates into a decisive play. Demonstrates mechanics and coordination viewers love to learn from.",
        criteria=dict(
            ui="Ability/ultimate icon state change (ready → activated/on-cooldown); ability HUD; full-screen ult VFX overlay.",
            nonui="Signature ability activation SFX / operator voice line; screen-space VFX (drone, smoke, healing aura).",
            trigger="Ability/ultimate activation followed by a payoff within ~8 s (≥2 kills or objective captured); optionally ≥2 different abilities chained = 'combo'.",
            thresholds="Activation + payoff (≥2 kills / objective) within 8–10 s; cooldown debounce to avoid re-trigger.",
            fp="Tie activation to a kill-feed payoff (don't clip every ability use); template-match specific ability icons; require VFX + kills to co-occur.",
        ),
        clip=dict(preroll="7 s (ability wind-up)", postroll="5 s", duration="20–35 s",
                  naming="DF_WAR_OperatorCombo_<operator>_<map>_<UTC>_<user>.mp4"),
    ),
    dict(
        idx=5, name="Clutch 1v4", mode="Both",
        justification="Peak skill expression — the last player alive winning while outnumbered. Classic viral, highly shareable clip material.",
        criteria=dict(
            ui="Teammate-status HUD shows all teammates down/dead (player = last alive); nearby enemy count; kill feed of player eliminating multiple enemies.",
            nonui="Drop in friendly comms/chatter; absence of revive prompts; tension music ramp.",
            trigger="All teammates dead (HUD) AND local player eliminates the remaining enemies while outnumbered (≥1v3 / 1v4) and survives/wins.",
            thresholds="Outnumbered ≥1v3 (configurable to 1v4); player secures remaining kills; engagement resolves in player's favour.",
            fp="Verify teammate-status HUD via template (all dead); count player eliminations from kill feed; require post-engagement survival.",
        ),
        clip=dict(preroll="12 s (to establish the odds)", postroll="6 s", duration="30–50 s",
                  naming="DF_<mode>_Clutch1v4_<map>_<UTC>_<user>.mp4"),
    ),
    dict(
        idx=6, name="Clutch in Overtime", mode="Warfare",
        justification="The highest-tension contextual moment — a match-deciding clutch during overtime. Maximum drama and narrative payoff.",
        criteria=dict(
            ui="'OVERTIME' banner/indicator on HUD/scoreboard; overtime objective/timer state; combined with clutch conditions (final objective, low tickets).",
            nonui="Overtime announcer voice line / audio sting; music intensity spike.",
            trigger="Overtime state active (OCR 'OVERTIME') AND a decisive player action (final kills / objective secure-defuse) leading to match end within the overtime window.",
            thresholds="Overtime flag = true; decisive action by player; match/round ends (win) shortly after the play.",
            fp="Require OVERTIME template AND a match-result/objective state change right after the play; trigger only on the impactful action, not every OT tick.",
        ),
        clip=dict(preroll="15 s", postroll="8 s (include the win/result moment)", duration="35–55 s",
                  naming="DF_WAR_ClutchOvertime_<map>_<UTC>_<user>.mp4"),
    ),
]


def clip_conventions_slide():
    s = new_slide()
    slide_title(s, "STEP 4 — CLIP MANAGEMENT", "Global Clip Conventions & Quality")

    # left: defaults table
    rect(s, Inches(0.6), Inches(1.55), Inches(6.0), Inches(5.25), fill=PANEL)
    _, tf = textbox(s, Inches(0.85), Inches(1.7), Inches(5.5), Inches(5.0))
    add_para(tf, "Pre/Post-roll Defaults", size=16, color=ACCENT, bold=True, first=True, space_after=8)
    table = [
        ("MandelBrick", "8 s / 5 s", "45–60 s"),
        ("Solo Squad Wipe", "10 s / 5 s", "25–40 s"),
        ("Tank Multi-Kill", "6 s / 4 s", "20–35 s"),
        ("Operator Combo", "7 s / 5 s", "20–35 s"),
        ("Clutch 1v4", "12 s / 6 s", "30–50 s"),
        ("Clutch Overtime", "15 s / 8 s", "35–55 s"),
    ]
    p = add_para(tf, "Event           Pre/Post        Length", size=12.5, color=ACCENT_2, bold=True, space_after=4)
    for name, pr, dur in table:
        p = tf.add_paragraph()
        p.space_after = Pt(5)
        set_run(p.add_run(), f"{name:<18}", size=12.5, color=TEXT, bold=True, font="Consolas")
        set_run(p.add_run(), f"{pr:<14}", size=12.5, color=MUTED, font="Consolas")
        set_run(p.add_run(), dur, size=12.5, color=GOOD, font="Consolas")

    # right: rules
    rect(s, Inches(6.8), Inches(1.55), Inches(5.93), Inches(5.25), fill=PANEL_2)
    _, tf = textbox(s, Inches(7.05), Inches(1.7), Inches(5.45), Inches(5.0))
    add_para(tf, "Management Rules", size=16, color=ACCENT, bold=True, first=True, space_after=8)
    rules = [
        "Naming: DF_<MODE>_<EVENT>_<context>_<UTC>_<user>.mp4 — sortable & self-describing.",
        "De-dup / cooldown: ≥20 s gap between clips; merge overlapping detection windows into one clip.",
        "Length guard: min 10 s, hard cap 60 s to keep clips snackable.",
        "Quality: match source resolution (1080p/1440p) at 30–60 fps, H.264 export.",
        "Metadata tags: mode, event type, map, kill count, confidence score.",
        "Confidence scoring: fuse UI + audio + kill-feed signals; auto-publish above threshold, queue the rest for review.",
    ]
    for r in rules:
        add_para(tf, r, size=13.5, color=TEXT, bullet=True, space_after=9)


def closing_slide():
    s = new_slide()
    rect(s, 0, 0, Inches(0.22), H, fill=ACCENT)
    slide_title(s, "SUMMARY", "Reliable Detection, Minimal False Positives")
    _, tf = textbox(s, Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.0))
    pts = [
        ("Sensor fusion is the key principle", "Each event requires ≥2 corroborating signals (UI + audio + scene/state) before a clip is cut — this is the main false-positive defence."),
        ("Six events span both modes", "Operations contributes objective- & skill-driven moments (MandelBrick, Solo Wipe); Warfare contributes spectacle & drama (Tank, Operator Combo, Clutch 1v4, Clutch Overtime)."),
        ("Context multiplies value", "Overtime and 1v4 conditions act as 'importance multipliers' — same play, far higher viewer payoff, so they get longer pre/post-roll."),
        ("Everything is tunable", "Thresholds, windows, and roll times are parameters — easy to A/B test against creator engagement to refine clip quality over time."),
    ]
    for head, body in pts:
        add_para(tf, head, size=17, color=ACCENT_2, bold=True, space_after=2, space_before=8)
        add_para(tf, body, size=14, color=TEXT, space_after=4)


def main():
    title_slide()
    game_selection_slide()
    methodology_slide()
    overview_slide()
    for ev in EVENTS:
        event_slide(**ev)
    clip_conventions_slide()
    closing_slide()
    out = "Task1_Game_Understanding_DeltaForce.pptx"
    prs.save(out)
    print(f"Saved {out} with {len(prs.slides._sldIdLst)} slides.")


if __name__ == "__main__":
    main()
